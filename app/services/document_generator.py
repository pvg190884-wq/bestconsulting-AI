"""Document Generator — создание файлов Excel, PowerPoint, PDF."""
import io
from app.utils.logger import setup_logging

logger = setup_logging()

FONT_PATH = "app/assets/fonts/DejaVuSans.ttf"


def build_xlsx(title: str, headers: list[str], rows: list[list]) -> bytes:
    """Создаёт Excel-файл (сводная таблица) из заголовков и строк данных."""
    import openpyxl
    from openpyxl.styles import Font, PatternFill

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = title[:31] if title else "Отчёт"

    header_fill = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF")

    for col_idx, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col_idx, value=h)
        cell.font = header_font
        cell.fill = header_fill

    for row_idx, row in enumerate(rows, 2):
        for col_idx, value in enumerate(row, 1):
            ws.cell(row=row_idx, column=col_idx, value=value)

    for col_idx, h in enumerate(headers, 1):
        ws.column_dimensions[chr(64 + col_idx)].width = max(15, len(str(h)) + 5)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def build_pptx(title: str, slides: list[dict]) -> bytes:
    """Создаёт презентацию PowerPoint. slides = [{"title": str, "bullets": [str, ...]}]"""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Bestconsulting AI Core"

    bullet_layout = prs.slide_layouts[1]
    for s in slides:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = s.get("title", "")
        body = slide.placeholders[1].text_frame
        bullets = s.get("bullets", [])
        if bullets:
            body.text = bullets[0]
            for b in bullets[1:]:
                p = body.add_paragraph()
                p.text = b
                p.level = 0

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_pdf(title: str, sections: list[dict]) -> bytes:
    """Создаёт PDF-доклад. sections = [{"heading": str, "text": str}]"""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()

    try:
        pdf.add_font("DejaVu", "", FONT_PATH, uni=True)
        pdf.add_font("DejaVu", "B", FONT_PATH, uni=True)
        font_name = "DejaVu"
    except Exception as e:
        logger.warning(f"[DocGen] Шрифт DejaVu не найден ({e}), используется Arial (без кириллицы)")
        font_name = "Arial"

    pdf.set_font(font_name, "B", 18)
    pdf.multi_cell(0, 10, title)
    pdf.ln(4)

    for sec in sections:
        heading = sec.get("heading", "")
        text = sec.get("text", "")
        if heading:
            pdf.set_font(font_name, "B", 14)
            pdf.multi_cell(0, 8, heading)
            pdf.ln(2)
        if text:
            pdf.set_font(font_name, "", 11)
            pdf.multi_cell(0, 6, text)
            pdf.ln(4)

    out = pdf.output(dest="S")
    if isinstance(out, str):
        out = out.encode("latin-1")
    return bytes(out)
