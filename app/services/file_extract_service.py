"""File Extract Service — извлечение текста из файлов (TXT, PDF, Excel, PowerPoint, JPG, PNG)."""
import io
import base64
from app.utils.logger import setup_logging

logger = setup_logging()

MAX_EXTRACT_CHARS = 15000
MIN_TEXT_LENGTH_FOR_VALID_PDF = 50  # если извлечено меньше — считаем PDF сканом
MAX_OCR_PAGES = 5  # ограничение страниц для OCR через vision-модель (контроль стоимости)


def _extract_txt(content: bytes) -> str:
    for enc in ("utf-8", "cp1251", "latin-1"):
        try:
            return content.decode(enc)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="ignore")


def _extract_pdf_text(content: bytes) -> str:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t)
    return "\n".join(text_parts)


def _extract_xlsx(content: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"=== Лист: {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs)
                    if t.strip():
                        slide_text.append(t)
        if slide_text:
            parts.append(f"=== Слайд {i} ===\n" + "\n".join(slide_text))
    return "\n".join(parts)


async def _extract_image_via_vision(content: bytes, llm_service) -> str:
    """Извлекает текст с изображения через vision-модель (gpt-4o-mini через OpenRouter)."""
    b64 = base64.b64encode(content).decode("utf-8")
    messages = [
        {
            "role": "user",
            "content": [
                {"type": "text", "text": "Извлеки и выведи ВЕСЬ читаемый текст с этого изображения дословно, максимально подробно, ничего не сокращай. Если текста нет — кратко опиши, что на изображении."},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}},
            ],
        }
    ]
    try:
        result = await llm_service.generate("openai", messages, temperature=0.0)
        extracted = result.get("content", "")
        logger.info(f"[FileExtract] Vision-извлечение: получено {len(extracted)} символов")
        return extracted
    except Exception as e:
        logger.error(f"[FileExtract] Ошибка vision-извлечения: {e}")
        return ""


async def _extract_pdf_via_ocr(content: bytes, llm_service) -> str:
    """Рендерит страницы PDF как изображения и распознаёт текст через vision-модель.
    Используется как fallback для сканированных PDF без текстового слоя."""
    try:
        import pypdfium2 as pdfium
    except Exception as e:
        logger.error(f"[FileExtract] pypdfium2 недоступен: {e}")
        return ""

    if llm_service is None:
        logger.warning("[FileExtract] OCR невозможен: llm_service не передан")
        return ""

    try:
        pdf = pdfium.PdfDocument(io.BytesIO(content))
        n_pages = min(len(pdf), MAX_OCR_PAGES)
        logger.info(f"[FileExtract] OCR: документ содержит {len(pdf)} стр., обрабатываем {n_pages}")
        parts = []
        for i in range(n_pages):
            page = pdf[i]
            bitmap = page.render(scale=2.0)
            pil_image = bitmap.to_pil()
            buf = io.BytesIO()
            pil_image.save(buf, format="JPEG", quality=85)
            page_bytes = buf.getvalue()
            logger.info(f"[FileExtract] OCR: страница {i + 1} отрендерена, {len(page_bytes)} байт")

            page_text = await _extract_image_via_vision(page_bytes, llm_service)
            logger.info(f"[FileExtract] OCR: страница {i + 1} — распознано {len(page_text)} символов")
            if page_text:
                parts.append(f"=== Страница {i + 1} ===\n{page_text}")

        result = "\n\n".join(parts)
        logger.info(f"[FileExtract] OCR: итого распознано {len(result)} символов")
        return result
    except Exception as e:
        logger.error(f"[FileExtract] Ошибка OCR-рендеринга PDF: {e}")
        return ""


async def extract_text(content: bytes, filename: str, llm_service=None) -> str:
    """Определяет тип файла по расширению и извлекает текст."""
    ext = (filename.rsplit(".", 1)[-1] if "." in filename else "").lower()
    try:
        if ext == "txt":
            text = _extract_txt(content)
        elif ext == "pdf":
            text = _extract_pdf_text(content)
            logger.info(f"[FileExtract] pdfplumber извлёк {len(text.strip())} символов из '{filename}'")
            if len(text.strip()) < MIN_TEXT_LENGTH_FOR_VALID_PDF:
                logger.info(f"[FileExtract] PDF '{filename}' похож на скан (мало текста) — пробуем OCR")
                ocr_text = await _extract_pdf_via_ocr(content, llm_service)
                if ocr_text and len(ocr_text.strip()) > len(text.strip()):
                    text = ocr_text
                else:
                    logger.warning(f"[FileExtract] OCR не дал результата лучше исходного ({len(ocr_text.strip())} символов) — используем что есть")
        elif ext in ("xlsx", "xls"):
            text = _extract_xlsx(content)
        elif ext == "pptx":
            text = _extract_pptx(content)
        elif ext in ("jpg", "jpeg", "png"):
            if llm_service is None:
                return ""
            text = await _extract_image_via_vision(content, llm_service)
        else:
            logger.warning(f"[FileExtract] Неподдерживаемый формат: {filename}")
            return ""
    except Exception as e:
        logger.error(f"[FileExtract] Ошибка извлечения из {filename}: {e}")
        return ""

    return text[:MAX_EXTRACT_CHARS] if text else ""
